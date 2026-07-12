# -*- coding: utf-8 -*-
"""Erstellt die allgemeine WebArs-Verkaufspräsentation als PowerPoint (16:9).
Basiert auf dem Design der Mustafa-Präsentation (A4-HTML), entpersonalisiert.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Farben (aus dem Original-CSS)
BG      = RGBColor(0xF0, 0xF0, 0xF1)
SURFACE = RGBColor(0xF8, 0xF8, 0xF9)
TEXT    = RGBColor(0x0F, 0x0E, 0x0D)
MUTED   = RGBColor(0x6B, 0x6B, 0x6E)
ACCENT  = RGBColor(0x18, 0x40, 0xFF)
ACC_LT  = RGBColor(0x78, 0x96, 0xFF)
DARK    = RGBColor(0x04, 0x07, 0x0F)
DARK2   = RGBColor(0x0A, 0x12, 0x40)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LN = RGBColor(0xD9, 0xD9, 0xDB)

F_DISPLAY = "Segoe UI Black"
F_BODY    = "Segoe UI"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rect(slide, x, y, w, h, color, rounded=False, line_color=None, line_w=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(kind, x, y, w, h)
    solid(sp, color)
    if rounded:
        try:
            sp.adjustments[0] = 0.045
        except Exception:
            pass
    if line_color is not None:
        sp.line.color.rgb = line_color
        sp.line.width = line_w or Pt(1)
    return sp


def tb(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
       space_after=Pt(0), line_spacing=None):
    """runs: Liste von Absätzen; jeder Absatz = Liste von (text, size, bold, color, font, italic)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        if line_spacing:
            p.line_spacing = line_spacing
        for (text, size, bold, color, font, italic) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
            r.font.italic = italic
    return box


def R(text, size, bold=False, color=TEXT, font=F_BODY, italic=False):
    return (text, size, bold, color, font, italic)


def label(slide, x, y, text, color=ACCENT):
    tb(slide, x, y, Inches(6), Inches(0.3),
       [[R(text.upper(), 11, True, color)]])


def footer(slide, page):
    ln = rect(slide, Inches(0.6), Inches(7.02), Inches(12.13), Pt(1), GREY_LN)
    tb(slide, Inches(0.6), Inches(7.1), Inches(4), Inches(0.3),
       [[R("WEBARS", 12, True, TEXT, F_DISPLAY)]])
    tb(slide, Inches(8.73), Inches(7.13), Inches(4), Inches(0.3),
       [[R(f"WEBSITE-KONZEPT — SEITE {page}", 9, False, MUTED)]], align=PP_ALIGN.RIGHT)


# ================= SEITE 1 — COVER =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARK)
# Blauer "Glow" unten links (Annäherung an das radiale Gradient-Cover)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-4.5), Inches(3.2), Inches(11), Inches(9))
solid(glow, DARK2)
glow2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3.5), Inches(5.2), Inches(7), Inches(6))
solid(glow2, RGBColor(0x11, 0x24, 0x8F))

tb(s, Inches(0.7), Inches(0.55), Inches(5), Inches(0.55),
   [[R("WEB", 26, True, WHITE, F_DISPLAY), R("ARS", 26, True, ACC_LT, F_DISPLAY)]])
tb(s, Inches(7.93), Inches(0.65), Inches(4.7), Inches(0.4),
   [[R("WEBSITE-KONZEPT 2026", 11, True, ACC_LT)]], align=PP_ALIGN.RIGHT)

rect(s, Inches(0.72), Inches(2.35), Inches(0.9), Pt(3), ACCENT)
tb(s, Inches(0.7), Inches(2.6), Inches(11), Inches(2.4),
   [[R("Ihre Website.", 52, True, WHITE, F_DISPLAY)],
    [R("Ihre neuen ", 52, True, WHITE, F_DISPLAY), R("Kunden.", 52, True, ACC_LT, F_DISPLAY, True)]],
   line_spacing=1.0)
tb(s, Inches(0.72), Inches(4.75), Inches(7.2), Inches(1),
   [[R("Ein maßgeschneidertes Angebot in drei Paketen — professionell umgesetzt, "
       "blitzschnell online und gebaut, um Anfragen zu bringen.", 14, False, RGBColor(0xC9, 0xCE, 0xE0))]],
   line_spacing=1.25)

tb(s, Inches(0.72), Inches(6.45), Inches(6), Inches(0.8),
   [[R("Erstellt für ", 12, False, RGBColor(0xA8, 0xAE, 0xC4)), R("[Kundenname]", 12, True, WHITE)],
    [R("[Ort], [Monat Jahr]", 12, False, RGBColor(0xA8, 0xAE, 0xC4))]], line_spacing=1.2)
tb(s, Inches(7.63), Inches(6.45), Inches(5), Inches(0.8),
   [[R("Teodor Turlea", 12, True, WHITE)],
    [R("WebArs e.U. — webars.at", 12, False, RGBColor(0xA8, 0xAE, 0xC4))]],
   align=PP_ALIGN.RIGHT, line_spacing=1.2)

# ================= SEITE 2 — WARUM JETZT =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, BG)
label(s, Inches(0.6), Inches(0.5), "Ausgangslage")
tb(s, Inches(0.6), Inches(0.82), Inches(10), Inches(0.6),
   [[R("Warum jetzt der richtige Zeitpunkt ist", 27, True, TEXT, F_DISPLAY)]])

cards = [
    ("Google Maps hat sich geändert",
     "Google bevorzugt Unternehmen mit eigener Website deutlich stärker in den lokalen Ergebnissen. "
     "Ohne Website verlieren Einträge sichtbar an Reichweite — die Konkurrenz rückt nach oben."),
    ("Kunden prüfen Sie online — vor dem Anruf",
     "Wer Sie empfohlen bekommt, sucht zuerst im Internet. Keine Website bedeutet für viele: "
     "der Auftrag geht an den Mitbewerber, der professionell auftritt."),
    ("Eine Website arbeitet rund um die Uhr",
     "Anfragen, Öffnungszeiten, Leistungen, Anfahrt — alles beantwortet sich von selbst. "
     "Sie gewinnen Kunden, während Sie arbeiten oder schlafen."),
]
y = Inches(1.72)
for title, body in cards:
    rect(s, Inches(0.6), y, Inches(7.35), Inches(1.55), SURFACE, rounded=True,
         line_color=GREY_LN, line_w=Pt(0.75))
    tb(s, Inches(0.85), y + Inches(0.18), Inches(6.9), Inches(1.25),
       [[R(title, 13.5, True, TEXT)],
        [R(body, 10.5, False, MUTED)]], line_spacing=1.15, space_after=Pt(4))
    y += Inches(1.73)

# Dunkles Panel rechts
rect(s, Inches(8.2), Inches(1.72), Inches(4.53), Inches(5.01), DARK, rounded=True)
tb(s, Inches(8.55), Inches(2.1), Inches(3.85), Inches(1.1),
   [[R("2 Tage", 46, True, ACC_LT, F_DISPLAY)]])
tb(s, Inches(8.55), Inches(3.2), Inches(3.85), Inches(1),
   [[R("Von der Entscheidung bis zur fertigen Website — das ist mein Versprechen.",
       13, True, WHITE)]], line_spacing=1.2)
rect(s, Inches(8.55), Inches(4.35), Inches(0.7), Pt(2), ACCENT)
tb(s, Inches(8.55), Inches(4.6), Inches(3.85), Inches(1.9),
   [[R("Kein monatelanges Projekt, keine Agentur-Warteschleife. Sie entscheiden, ich setze um — "
       "und in zwei Tagen sind Sie online, professionell und auffindbar.",
       11, False, RGBColor(0xB8, 0xBE, 0xD2))]], line_spacing=1.25)
footer(s, 2)

# ================= SEITE 3 — PAKETE =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, BG)
label(s, Inches(0.6), Inches(0.38), "Ihre Möglichkeiten")
tb(s, Inches(0.6), Inches(0.68), Inches(7), Inches(0.55),
   [[R("Drei Pakete — ein Ziel", 26, True, TEXT, F_DISPLAY)]])
tb(s, Inches(8.1), Inches(0.62), Inches(4.6), Inches(0.7),
   [[R("Alle Preise einmalig, ohne versteckte Kosten. Hosting und Domain im ersten Jahr inklusive.",
       9.5, False, MUTED)]], align=PP_ALIGN.RIGHT, line_spacing=1.15)

packs = [
    ("Start", "500 €", "Der solide Einstieg: eine professionelle Visitenkarte im Internet.",
     ["Moderne Onepage-Website (1 Seite)", "Optimiert für Handy & Tablet",
      "Kontaktformular & Anruf-Button", "Verknüpfung mit Google Maps",
      "Impressum & Datenschutz inklusive", "1 Korrekturrunde"],
     "Fertig in 2 Tagen", False),
    ("Business", "1.000 €", "Der komplette Marktauftritt: mehrere Seiten, gefunden werden, Anfragen gewinnen.",
     ["Bis zu 5 individuelle Unterseiten", "Maßgeschneidertes Design — kein Baukasten",
      "Google-Optimierung (SEO) — Sie werden gefunden", "Volle Google-Maps-Anbindung",
      "WhatsApp-Direktkontakt für Kunden", "Professionelle Texte — ich schreibe für Sie",
      "2 Korrekturrunden", "3 Monate persönlicher Support inklusive"],
     "Fertig in 2 Tagen — Empfehlung", True),
    ("Premium", "1.500 €", "Für den maximalen Auftritt mit Ausbau-Reserven.",
     ["Alles aus Business, zusätzlich:", "Bis zu 10 Unterseiten", "News- / Aktuelles-Bereich",
      "Zweite Sprache (z. B. Englisch)", "Beratung zu Fotos & Außenauftritt",
      "6 Monate persönlicher Support", "Bevorzugte Bearbeitung bei Änderungen"],
     "Fertig in 3–4 Tagen", False),
]
x = Inches(0.6)
cw, gap = Inches(3.98), Inches(0.29)
for name, price, desc, items, foot, featured in packs:
    cy, ch = Inches(1.55), Inches(5.32)
    if featured:
        rect(s, x, Inches(1.32), cw, Inches(0.34), ACCENT, rounded=True)
        tb(s, x, Inches(1.36), cw, Inches(0.28),
           [[R("MEISTGEWÄHLT — BESTER GEGENWERT", 9, True, WHITE)]], align=PP_ALIGN.CENTER)
        rect(s, x, cy, cw, ch, SURFACE, rounded=True, line_color=ACCENT, line_w=Pt(2))
    else:
        rect(s, x, cy, cw, ch, SURFACE, rounded=True, line_color=GREY_LN, line_w=Pt(0.75))
    ix = x + Inches(0.28)
    iw = cw - Inches(0.56)
    tb(s, ix, cy + Inches(0.22), iw, Inches(0.35),
       [[R(name.upper(), 13, True, ACCENT if featured else TEXT)]])
    tb(s, ix, cy + Inches(0.55), iw, Inches(0.55),
       [[R(price, 27, True, TEXT, F_DISPLAY), R("  einmalig", 10, False, MUTED)]])
    tb(s, ix, cy + Inches(1.18), iw, Inches(0.6),
       [[R(desc, 10, False, MUTED)]], line_spacing=1.12)
    rows = [[R("•  ", 10.5, True, ACCENT), R(it, 10.5, False, TEXT)] for it in items]
    tb(s, ix, cy + Inches(1.85), iw, Inches(2.75), rows, line_spacing=1.06, space_after=Pt(3.5))
    rect(s, ix, cy + ch - Inches(0.62), iw, Pt(1), GREY_LN)
    tb(s, ix, cy + ch - Inches(0.5), iw, Inches(0.35),
       [[R(foot, 9.5, True, ACCENT if featured else MUTED)]])
    x += cw + gap
footer(s, 3)

# ================= SEITE 4 — ABLAUF =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, BG)
label(s, Inches(0.6), Inches(0.5), "So einfach geht es")
tb(s, Inches(0.6), Inches(0.82), Inches(11), Inches(0.6),
   [[R("Von heute bis online — in 4 Schritten", 27, True, TEXT, F_DISPLAY)]])

steps = [
    ("01", "Gespräch & Entscheidung",
     "Wir klären gemeinsam, was Ihre Website leisten soll: Leistungen, Zielkunden, gewünschte Seiten. Sie wählen Ihr Paket.", "HEUTE"),
    ("02", "Design-Entwurf",
     "Ich gestalte Ihren Auftritt: Farben, Aufbau, Struktur — abgestimmt auf Ihr Unternehmen. Sie sehen ein konkretes Ergebnis, kein Konzeptpapier.", "TAG 1"),
    ("03", "Umsetzung & Feinschliff",
     "Die komplette Website wird gebaut: alle Seiten, Texte, Kontaktwege, Google-Anbindung. Ihre Korrekturwünsche fließen direkt ein.", "TAG 1–2"),
    ("04", "Livegang & Übergabe",
     "Ihre Website geht online — mit eigener Adresse, auf Google auffindbar, auf jedem Gerät perfekt. Alles wird übergeben und in Ruhe erklärt.", "TAG 2"),
]
x = Inches(0.6)
cw, gap = Inches(2.95), Inches(0.29)
for num, title, body, when in steps:
    rect(s, x, Inches(1.75), cw, Inches(3.55), SURFACE, rounded=True,
         line_color=GREY_LN, line_w=Pt(0.75))
    ix, iw = x + Inches(0.24), cw - Inches(0.48)
    tb(s, ix, Inches(1.98), iw, Inches(0.5),
       [[R(num, 24, True, ACCENT, F_DISPLAY)]])
    tb(s, ix, Inches(2.55), iw, Inches(0.65),
       [[R(title, 13, True, TEXT)]], line_spacing=1.05)
    tb(s, ix, Inches(3.18), iw, Inches(1.7),
       [[R(body, 10, False, MUTED)]], line_spacing=1.18)
    tb(s, ix, Inches(4.88), iw, Inches(0.3),
       [[R(when, 9.5, True, ACCENT)]])
    x += cw + gap

rect(s, Inches(0.6), Inches(5.6), Inches(12.13), Inches(1.1), DARK, rounded=True)
tb(s, Inches(0.95), Inches(5.78), Inches(11.4), Inches(0.75),
   [[R("Ihr Aufwand: rund eine Stunde. ", 12, True, WHITE),
     R("Sie liefern mir Ihr Logo, ein paar Fotos und Ihre Eckdaten — den Rest übernehme ich komplett. "
       "Texte, Technik, Design, Google: alles aus einer Hand.", 12, False, RGBColor(0xB8, 0xBE, 0xD2))]],
   line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# ================= SEITE 5 — NÄCHSTER SCHRITT =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARK)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-4), Inches(11), Inches(9))
solid(glow, DARK2)

label(s, Inches(0.7), Inches(0.75), "Der nächste Schritt", color=ACC_LT)
tb(s, Inches(0.7), Inches(1.1), Inches(11.5), Inches(1.5),
   [[R("Entscheiden Sie heute — ", 36, True, WHITE, F_DISPLAY),
     R("online in 2 Tagen.", 36, True, ACC_LT, F_DISPLAY, True)]], line_spacing=1.05)
tb(s, Inches(0.72), Inches(2.35), Inches(8.5), Inches(0.9),
   [[R("Sie wählen das Paket, das zu Ihnen passt. Ich starte sofort mit der Umsetzung — "
       "und noch diese Woche präsentiert sich Ihr Unternehmen professionell im Internet.",
       13.5, False, RGBColor(0xC9, 0xCE, 0xE0))]], line_spacing=1.3)

boxes = [
    ("FESTPREIS", "Keine versteckten Kosten, keine Abos — Sie wissen vorher genau, was es kostet."),
    ("PERSÖNLICH", "Ein direkter Ansprechpartner — kein Callcenter, keine Agentur-Warteschleife."),
    ("AUS WIEN", "Kurze Wege, schnelle Abstimmung, Handschlagqualität — persönlich statt anonym."),
]
x = Inches(0.7)
cw, gap = Inches(3.9), Inches(0.26)
for k, v in boxes:
    bx = rect(s, x, Inches(3.6), cw, Inches(1.85), RGBColor(0x0C, 0x11, 0x26), rounded=True,
              line_color=RGBColor(0x24, 0x2E, 0x55), line_w=Pt(0.75))
    tb(s, x + Inches(0.26), Inches(3.85), cw - Inches(0.52), Inches(0.3),
       [[R(k, 11, True, ACC_LT)]])
    tb(s, x + Inches(0.26), Inches(4.22), cw - Inches(0.52), Inches(1.1),
       [[R(v, 10.5, False, RGBColor(0xC9, 0xCE, 0xE0))]], line_spacing=1.2)
    x += cw + gap

rect(s, Inches(0.7), Inches(6.15), Inches(11.93), Pt(1), RGBColor(0x24, 0x2E, 0x55))
tb(s, Inches(0.7), Inches(6.35), Inches(8), Inches(0.8),
   [[R("Teodor Turlea", 13, True, WHITE), R(" — WebArs e.U.", 13, False, RGBColor(0xC9, 0xCE, 0xE0))],
    [R("turlea@webars.at   ·   webars.at", 11.5, False, RGBColor(0xA8, 0xAE, 0xC4))]], line_spacing=1.25)
tb(s, Inches(8.63), Inches(6.5), Inches(4), Inches(0.5),
   [[R("WEB", 22, True, WHITE, F_DISPLAY), R("ARS", 22, True, ACC_LT, F_DISPLAY)]],
   align=PP_ALIGN.RIGHT)

out = r"C:\Users\Drsacula\Claude-Dateien\webars-praesentation-vorlage\WebArs-Praesentation-Vorlage.pptx"
prs.save(out)
print("OK:", out)
